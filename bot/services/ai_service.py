import io
import json
import logging
import asyncio
import urllib.request
import urllib.parse
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Pt as DocxPt, Cm as DocxCm, Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_ORIENT
import qrcode

from config import (
    AI_PROVIDER, GEMINI_API_KEY, GEMINI_API_KEYS, GEMINI_MODELS,
    OPENAI_API_KEY, OPENAI_MODEL, PEXELS_API_KEY, PIXABAY_API_KEY
)

logger = logging.getLogger(__name__)

# ===== GEMINI SETUP =====
genai.configure(api_key=GEMINI_API_KEY)


# ============================================================
# UNIVERSAL AI CALL
# ============================================================

async def ai_generate(prompt: str, max_tokens: int = 4000, temperature: float = 0.7) -> str:
    """AI orqali matn generatsiya qilish."""
    if AI_PROVIDER == "openai":
        return await _openai_generate(prompt, max_tokens, temperature)
    return await _gemini_generate(prompt, max_tokens, temperature)


async def _gemini_generate(prompt: str, max_tokens: int = 4000, temperature: float = 0.7) -> str:
    """
    Gemini — barcha API keylarni va modellarni ketma-ket sinaydi.
    Key1 + barcha modellar → Key2 + barcha modellar → ... → Kutish → Qayta sinash
    """
    last_error = None
    max_retries = 3
    
    for attempt in range(max_retries):
        # Har bir API key bilan sinash
        for key_idx, api_key in enumerate(GEMINI_API_KEYS):
            genai.configure(api_key=api_key)
            
            for model_name in GEMINI_MODELS:
                try:
                    logger.info(f"Trying model: {model_name}, key #{key_idx + 1} (attempt {attempt + 1})")
                    model = genai.GenerativeModel(model_name)
                    generation_config = genai.types.GenerationConfig(
                        max_output_tokens=max_tokens,
                        temperature=temperature,
                    )
                    response = await model.generate_content_async(
                        prompt, generation_config=generation_config,
                    )
                    if response and response.text:
                        logger.info(f"✅ Success: {model_name}, key #{key_idx + 1}")
                        return response.text.strip()
                    else:
                        continue
                except Exception as e:
                    last_error = e
                    error_msg = str(e)
                    if "429" in error_msg or "ResourceExhausted" in error_msg or "quota" in error_msg.lower():
                        logger.warning(f"⚠️ Quota: {model_name}, key #{key_idx + 1}")
                    elif "404" in error_msg or "NotFound" in error_msg:
                        logger.warning(f"❌ Not found: {model_name}")
                    else:
                        logger.warning(f"❌ Error: {model_name}: {e}")
                    continue
        
        # Barcha keylar va modellar ishlamadi — kutish
        if attempt < max_retries - 1:
            wait_time = 20 * (attempt + 1)  # 20s, 40s, 60s
            logger.info(f"⏳ All keys/models failed. Waiting {wait_time}s...")
            await asyncio.sleep(wait_time)
    
    raise Exception(f"❌ Barcha API keylar va modellar ishlamadi. Oxirgi xato: {last_error}")


async def _openai_generate(prompt: str, max_tokens: int = 4000, temperature: float = 0.7) -> str:
    """OpenAI orqali generatsiya (O'CHIRILGAN)."""
    try:
        import openai
        client = openai.AsyncOpenAI(api_key=OPENAI_API_KEY)
        response = await client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=temperature,
            max_tokens=max_tokens,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        raise Exception(f"OpenAI xatosi: {e}")


# ============================================================
# GOST STANDARTLARI
# ============================================================

# GOST: Referat/Mustaqil ish sahifa sozlamalari
GOST_DOC = {
    "font_name": "Times New Roman",
    "font_size": DocxPt(14),
    "line_spacing": 1.5,
    "align": WD_ALIGN_PARAGRAPH.JUSTIFY,
    "first_line_indent": DocxCm(1.25),
    "margin_left": DocxCm(3.0),
    "margin_right": DocxCm(1.5),
    "margin_top": DocxCm(2.0),
    "margin_bottom": DocxCm(2.0),
}

# GOST: Slayd shriftlari
GOST_PPT = {
    "title_font": "Times New Roman",
    "title_size": Pt(32),      # 28-36 pt
    "body_font": "Times New Roman",
    "body_size": Pt(20),       # 18-24 pt
}


def _apply_gost_paragraph(paragraph, font_name="Times New Roman", font_size=None,
                           bold=False, alignment=None, first_indent=None):
    """GOST formatini paragrafga qo'llash."""
    pf = paragraph.paragraph_format
    if alignment:
        pf.alignment = alignment
    if first_indent:
        pf.first_line_indent = first_indent
    pf.line_spacing = 1.5
    
    for run in paragraph.runs:
        run.font.name = font_name
        if font_size:
            run.font.size = font_size
        run.font.bold = bold


def _setup_gost_document(doc: Document):
    """GOST sahifa sozlamalarini o'rnatish: A4, maydonlar, shrift."""
    for section in doc.sections:
        section.page_width = DocxCm(21.0)   # A4
        section.page_height = DocxCm(29.7)  # A4
        section.left_margin = GOST_DOC["margin_left"]
        section.right_margin = GOST_DOC["margin_right"]
        section.top_margin = GOST_DOC["margin_top"]
        section.bottom_margin = GOST_DOC["margin_bottom"]


def _add_gost_title_page(doc: Document, title: str, doc_type: str, lang: str):
    """GOST titul varag'ini yaratish."""
    # Universitet nomi
    uni_text = {
        "uz": "O'ZBEKISTON RESPUBLIKASI\nOLIY VA O'RTA MAXSUS TA'LIM VAZIRLIGI",
        "ru": "РЕСПУБЛИКА УЗБЕКИСТАН\nМИНИСТЕРСТВО ВЫСШЕГО И СРЕДНЕГО СПЕЦИАЛЬНОГО ОБРАЗОВАНИЯ",
        "en": "REPUBLIC OF UZBEKISTAN\nMINISTRY OF HIGHER AND SECONDARY SPECIALIZED EDUCATION",
    }
    
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(uni_text.get(lang, uni_text["uz"]))
    run.font.name = "Times New Roman"
    run.font.size = DocxPt(14)
    
    # Bo'sh joy
    for _ in range(4):
        doc.add_paragraph()
    
    # Hujjat turi
    type_names = {
        "uz": {"referat": "REFERAT", "mustaqil_ish": "MUSTAQIL ISH"},
        "ru": {"referat": "РЕФЕРАТ", "mustaqil_ish": "САМОСТОЯТЕЛЬНАЯ РАБОТА"},
        "en": {"referat": "ESSAY/REPORT", "mustaqil_ish": "INDEPENDENT WORK"},
    }
    type_name = type_names.get(lang, type_names["uz"]).get(doc_type, "REFERAT")
    
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(type_name)
    run.font.name = "Times New Roman"
    run.font.size = DocxPt(20)
    run.font.bold = True
    
    # Bo'sh joy
    doc.add_paragraph()
    
    # Mavzu
    mavzu_label = {"uz": "Mavzu:", "ru": "Тема:", "en": "Topic:"}
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"{mavzu_label.get(lang, 'Mavzu:')} {title}")
    run.font.name = "Times New Roman"
    run.font.size = DocxPt(16)
    run.font.bold = True
    
    # Bo'sh joy
    for _ in range(8):
        doc.add_paragraph()
    
    # Sana
    import time
    from datetime import datetime
    now = datetime.now()
    city = {"uz": "Toshkent", "ru": "Ташкент", "en": "Tashkent"}
    
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"{city.get(lang, 'Toshkent')} — {now.year}")
    run.font.name = "Times New Roman"
    run.font.size = DocxPt(14)
    
    doc.add_page_break()


def _add_gost_heading(doc: Document, text: str, level: int = 1):
    """GOST sarlavha qo'shish: Bold, markazda, oxirida nuqta yo'q."""
    p = doc.add_paragraph()
    if level == 1:
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    else:
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    
    # Oxirida nuqtani olib tashlash
    clean_text = text.rstrip(".")
    
    run = p.add_run(clean_text.upper() if level == 1 else clean_text)
    run.font.name = "Times New Roman"
    run.font.size = DocxPt(14)
    run.font.bold = True
    
    # Bo'sh joy
    p.paragraph_format.space_before = DocxPt(12)
    p.paragraph_format.space_after = DocxPt(6)
    
    return p


def _add_gost_paragraph(doc: Document, text: str):
    """GOST asosiy matn paragrafi: 14pt, 1.5 interval, abzats 1.25sm, Justify."""
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.first_line_indent = DocxCm(1.25)
    p.paragraph_format.line_spacing = 1.5
    
    run = p.add_run(text)
    run.font.name = "Times New Roman"
    run.font.size = DocxPt(14)
    
    return p


# ============================================================
# PPT GENERATION — GOST
# ============================================================

PPT_COLOR_SCHEMES = {
    "business": {"bg": "1F3864", "title": "FFFFFF", "text": "D9E2F3"},
    "minimal": {"bg": "FFFFFF", "title": "2D2D2D", "text": "555555"},
    "dark": {"bg": "1A1A2E", "title": "E94560", "text": "EAEAEA"},
    "modern": {"bg": "F8F9FA", "title": "212529", "text": "495057"},
    "education": {"bg": "E8F5E9", "title": "1B5E20", "text": "2E7D32"},
    "corporate": {"bg": "0D47A1", "title": "FFFFFF", "text": "BBDEFB"},
    "startup": {"bg": "FFF3E0", "title": "E65100", "text": "F57C00"},
    "creative": {"bg": "F3E5F5", "title": "6A1B9A", "text": "8E24AA"},
    "elegant": {"bg": "FBE9E7", "title": "3E2723", "text": "5D4037"},
    "premium": {"bg": "212121", "title": "FFD700", "text": "BDBDBD"},
    "ocean": {"bg": "01579B", "title": "FFFFFF", "text": "B3E5FC"},
    "sunset": {"bg": "BF360C", "title": "FFF3E0", "text": "FFCCBC"},
    "forest": {"bg": "1B5E20", "title": "FFFFFF", "text": "C8E6C9"},
    "royal": {"bg": "4A148C", "title": "FFD700", "text": "E1BEE7"},
    "neon": {"bg": "0D0D0D", "title": "00E5FF", "text": "B2EBF2"},
}

# GOST slayd tuzilishi
GOST_SLIDE_STRUCTURE = {
    "uz": [
        "Titul slaydi",
        "Mavzu dolzarbligi",
        "Maqsad va vazifalar",
        # ... asosiy qism ...
        "Natijalar",
        "Xulosa",
        "Foydalanilgan adabiyotlar",
    ],
    "ru": [
        "Титульный слайд",
        "Актуальность темы",
        "Цель и задачи",
        "Результаты",
        "Заключение",
        "Использованная литература",
    ],
    "en": [
        "Title Slide",
        "Topic Relevance",
        "Goals and Objectives",
        "Results",
        "Conclusion",
        "References",
    ],
}


def generate_designs_preview() -> io.BytesIO:
    """15 ta dizaynni mini-slayd ko'rinishida bitta rasmga chizadi (3x5 grid)."""
    from PIL import Image, ImageDraw, ImageFont

    designs = [
        ("Business", "business"), ("Minimal", "minimal"), ("Dark", "dark"),
        ("Modern", "modern"), ("Education", "education"), ("Corporate", "corporate"),
        ("Startup", "startup"), ("Creative", "creative"), ("Elegant", "elegant"),
        ("Premium", "premium"), ("Ocean", "ocean"), ("Sunset", "sunset"),
        ("Forest", "forest"), ("Royal", "royal"), ("Neon", "neon"),
    ]
    cols, rows = 3, 5
    cw, ch = 340, 215
    pad = 18
    W = cols * cw + (cols + 1) * pad
    H = rows * ch + (rows + 1) * pad
    img = Image.new("RGB", (W, H), (13, 17, 28))
    d = ImageDraw.Draw(img)

    def load_font(size, bold=False):
        names = (["DejaVuSans-Bold.ttf", "arialbd.ttf", "Arial Bold.ttf", "arial.ttf"]
                 if bold else ["DejaVuSans.ttf", "arial.ttf", "Arial.ttf"])
        for n in names:
            try:
                return ImageFont.truetype(n, size)
            except Exception:
                continue
        return ImageFont.load_default()

    f_title = load_font(22, bold=True)
    f_label = load_font(15)

    def hx(h):
        return tuple(int(h[j:j+2], 16) for j in (0, 2, 4))

    for i, (name, key) in enumerate(designs):
        c = PPT_COLOR_SCHEMES.get(key, PPT_COLOR_SCHEMES["business"])
        r = i // cols
        col = i % cols
        x = pad + col * (cw + pad)
        y = pad + r * (ch + pad)
        sh = ch - 38
        bg, title, text = hx(c["bg"]), hx(c["title"]), hx(c["text"])

        # slayd foni
        d.rectangle([x, y, x + cw, y + sh], fill=bg)
        # accent burchak
        d.ellipse([x + cw - 46, y - 14, x + cw + 14, y + 46], fill=title)
        # sarlavha
        d.text((x + 18, y + 20), name, fill=title, font=f_title)
        d.rectangle([x + 18, y + 56, x + 18 + 130, y + 61], fill=title)
        # matn chiziqlari
        for k in range(3):
            ly = y + 82 + k * 24
            w = (cw - 60) if k < 2 else (cw - 60) // 2
            d.rectangle([x + 18, ly, x + 18 + w, ly + 7], fill=text)
        # pastki yorliq (raqam + nom)
        d.text((x + 18, y + sh + 9), f"#{i + 1}  {name}", fill=(205, 210, 220), font=f_label)

    out = io.BytesIO()
    img.save(out, format="PNG")
    out.seek(0)
    return out


async def generate_ppt_content(topic: str, slides: int, purpose: str, lang: str, extra: str = "", is_pro: bool = False) -> list:
    """
    GOST standartidagi taqdimot — avval REJA tuziladi, keyin har slayd real sarlavha bilan.
    Reja nomlari mazmunli bo'ladi, "1-bo'lim" emas.
    """
    lang_name = {"uz": "O'zbek", "ru": "Русский", "en": "English"}.get(lang, "O'zbek")

    # Asosiy qism slaydlari soni = jami - (titul + reja + kirish + xulosa + adabiyotlar)
    main_count = max(1, slides - 5)

    # 1-QADAM: REJA tuzish — mazmunli sarlavhalar
    plan_prompt = f"""Sen "{lang_name}" tilida GOST standartidagi akademik taqdimot (PPT) tayyorlayotgan mutaxassissan.
Mavzu: "{topic}"
Maqsad: {purpose}

VAZIFA: Asosiy qism uchun aniq {main_count} ta MAZMUNLI slayd sarlavhasini yoz.

QOIDALAR:
- Har bir sarlavha mavzuning aniq jihatini ifodalasin
- "1-bo'lim", "2-slayd", "Asosiy qism" kabi UMUMIY nomlar TAQIQLANADI
- Sarlavhalar qisqa, aniq va mavzuga oid bo'lsin (3-7 so'z)

Faqat JSON massiv qaytar:
["birinchi mazmunli sarlavha", "ikkinchi mazmunli sarlavha", ...]
Faqat to'g'ri JSON. Markdown yo'q."""

    main_titles = []
    try:
        pt = await ai_generate(plan_prompt, max_tokens=min(3000, 600 + main_count * 60), temperature=0.7)
        pt = pt.strip()
        if pt.startswith("```"):
            pt = pt.split("\n", 1)[1]
            if pt.endswith("```"):
                pt = pt[:-3]
        parsed = json.loads(pt.strip())
        if isinstance(parsed, list):
            main_titles = [str(t).strip() for t in parsed if str(t).strip()]
    except Exception:
        main_titles = []

    # Validatsiya
    def _bad(t):
        tl = (t or "").lower().strip()
        return (not tl) or len(tl) < 6 or tl.startswith(("1-", "2-", "3-", "slayd", "bob", "bo'lim"))

    # Yomon sarlavhalarni olib tashlaymiz (hammasini emas — faqat yomonlarini)
    main_titles = [t for t in main_titles if not _bad(t)]

    # Agar yetarli bo'lmasa, mavzuga oid jihatlar bilan to'ldiramiz (takrorsiz)
    if len(main_titles) < main_count:
        aspects = [
            "mohiyati va asosiy tushunchalari", "tarixiy rivojlanishi", "zamonaviy holati",
            "asosiy tamoyillari", "turlari va tasnifi", "ahamiyati va roli",
            "amaliy qo'llanilishi", "qiyosiy tahlili", "muammolari va yechimlari",
            "xalqaro tajriba", "statistik ko'rsatkichlari", "rivojlanish istiqbollari",
            "samaradorligini oshirish yo'llari", "ta'siri va natijalari", "kelajak tendensiyalari",
        ]
        i = 0
        while len(main_titles) < main_count:
            main_titles.append(f"{topic}: {aspects[i % len(aspects)]}")
            i += 1
    main_titles = main_titles[:main_count]

    # GOST slayd tuzilishi: titul + reja + kirish + [asosiy] + xulosa + adabiyotlar
    intro_t = {"uz": "Kirish", "ru": "Введение", "en": "Introduction"}.get(lang, "Kirish")
    plan_t = {"uz": "Reja", "ru": "План", "en": "Plan"}.get(lang, "Reja")
    concl_t = {"uz": "Xulosa", "ru": "Заключение", "en": "Conclusion"}.get(lang, "Xulosa")
    refs_t = {"uz": "Foydalanilgan adabiyotlar", "ru": "Использованная литература", "en": "References"}.get(lang, "Foydalanilgan adabiyotlar")

    # To'liq slayd sarlavhalari ro'yxati
    full_titles = [topic, plan_t, intro_t] + main_titles + [concl_t, refs_t]

    # 2-QADAM: Kontentni KICHIK GURUHLARGA bo'lib generatsiya qilamiz.
    # MUHIM: bitta katta so'rovda barcha slaydlarni so'rasak, AI JSON ni
    # yarmida kesib qo'yadi -> slaydlar yetishmaydi (masalan 30 o'rniga 15).
    # Shuning uchun har safar 5 tadan slayd so'raymiz va natijalarni birlashtiramiz.
    bullets_per_slide = "6-8" if is_pro else "5-6"
    pro_note = "\n- PRO sifat: chuqurroq tahlil, aniq raqamlar, statistika va misollar bilan boyit" if is_pro else ""

    async def _gen_batch(batch_titles: list) -> list:
        titles_block = "\n".join([f'{j+1}. "{t}"' for j, t in enumerate(batch_titles)])
        prompt = f"""Sen "{lang_name}" tilida GOST standartidagi akademik taqdimot kontentini yozayotgan mutaxassissan.
Mavzu: "{topic}"
{f'Qo`shimcha talab: {extra}' if extra else ''}

Quyidagi {len(batch_titles)} ta slayd uchun kontent yoz. Sarlavhalarni AYNAN shu holida ishlat, O'ZGARTIRMA:
{titles_block}

QOIDALAR:
- Agar sarlavha "Reja"/"План"/"Plan" bo'lsa — taqdimotning asosiy bo'limlari ro'yxatini bullet qil
- Agar sarlavha "Foydalanilgan adabiyotlar"/"References" bo'lsa — 6-7 ta real manba (GOST formatda)
- Boshqa slaydlarda {bullets_per_slide} ta to'liq, ma'lumotli bullet (har biri 12-22 so'z, aniq faktlar bilan)
- Bulletlar bir-birini takrorlamasin, akademik va mazmunli bo'lsin{pro_note}

Aniq {len(batch_titles)} ta obyektli JSON massiv qaytar. Har obyekt:
{{"title": "<yuqoridagi aynan o'sha sarlavha>", "content": ["bullet1", "bullet2", ...], "notes": "ma'ruzachi uchun 2 jumla"}}

Faqat to'g'ri, to'liq JSON massiv. Markdown yo'q."""
        for attempt in range(3):
            try:
                raw = await ai_generate(prompt, max_tokens=min(8000, 1500 + len(batch_titles) * 900), temperature=0.7)
                raw = raw.strip()
                if raw.startswith("```"):
                    raw = raw.split("\n", 1)[1]
                    if raw.endswith("```"):
                        raw = raw[:-3]
                arr = json.loads(raw.strip())
                if isinstance(arr, list) and arr:
                    return arr
            except Exception as e:
                logger.warning(f"PPT batch attempt {attempt + 1} failed: {e}")
        return []

    results = []
    BATCH = 5
    for bi in range(0, len(full_titles), BATCH):
        batch_titles = full_titles[bi:bi + BATCH]
        arr = await _gen_batch(batch_titles)
        # Har slaydni sarlavhasiga majburan bog'laymiz va yetishmaganini to'ldiramiz
        for k, t in enumerate(batch_titles):
            if k < len(arr) and isinstance(arr[k], dict) and arr[k].get("content"):
                item = arr[k]
                item["title"] = t
                results.append(item)
            else:
                # Zaxira — slayd hech qachon bo'sh qolmasin
                results.append({
                    "title": t,
                    "content": [f"{t} bo'yicha asosiy ma'lumotlar va tahlil",
                                "Mavzuning muhim jihatlari va xulosalari"],
                    "notes": ""
                })

    # "Reja" slaydini kafolatlangan tarzda bo'limlar ro'yxati bilan to'ldiramiz
    for item in results:
        if item.get("title") == plan_t:
            item["content"] = main_titles[:]
            break

    return results[:len(full_titles)]


def _has_image_source() -> bool:
    """Rasm manbasi doimo mavjud — Openverse kalitsiz ishlaydi."""
    return True


_UA = "Mozilla/5.0 (MasterStudentBot)"


def _download_bytes(url: str, headers: dict = None, timeout: int = 12) -> io.BytesIO:
    req = urllib.request.Request(url, headers=headers or {"User-Agent": _UA})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return io.BytesIO(resp.read())


def _fetch_image(query: str) -> io.BytesIO:
    """
    Mavzuga oid rasm yuklab oladi.
    Tartib: Pexels (kalit) -> Pixabay (kalit) -> Openverse (KALITSIZ, doim ishlaydi).
    """
    import random as _r

    # 1) PEXELS (eng sifatli, kalit kerak)
    if PEXELS_API_KEY:
        try:
            url = f"https://api.pexels.com/v1/search?query={urllib.parse.quote(query)}&per_page=8&orientation=landscape"
            data = json.loads(_download_bytes(url, {"Authorization": PEXELS_API_KEY}).read().decode())
            photos = data.get("photos", [])
            if photos:
                return _download_bytes(_r.choice(photos)["src"]["large"])
        except Exception as e:
            logger.warning(f"Pexels ishlamadi: {e}")

    # 2) PIXABAY (kalit kerak)
    if PIXABAY_API_KEY:
        try:
            url = (f"https://pixabay.com/api/?key={PIXABAY_API_KEY}"
                   f"&q={urllib.parse.quote(query)}&image_type=photo&orientation=horizontal&per_page=8&safesearch=true")
            data = json.loads(_download_bytes(url).read().decode())
            hits = data.get("hits", [])
            if hits:
                h = _r.choice(hits)
                return _download_bytes(h.get("largeImageURL") or h.get("webformatURL"))
        except Exception as e:
            logger.warning(f"Pixabay ishlamadi: {e}")

    # 3) OPENVERSE — KALITSIZ, hamma uchun ishlaydi (asosiy zaxira)
    try:
        url = f"https://api.openverse.org/v1/images/?q={urllib.parse.quote(query)}&page_size=10&mature=false"
        data = json.loads(_download_bytes(url, {"User-Agent": _UA}, timeout=15).read().decode())
        results = [r for r in data.get("results", []) if r.get("url")]
        if results:
            for r in _r.sample(results, min(3, len(results))):
                try:
                    return _download_bytes(r["url"], {"User-Agent": _UA})
                except Exception:
                    continue
    except Exception as e:
        logger.warning(f"Openverse ishlamadi: {e}")

    return None


async def _get_image_keyword(topic: str, slide_title: str, purpose: str = "") -> str:
    """AI yordamida slayd uchun mos rasm kalit so'zini (inglizcha) oladi. Maqsadga moslashtiriladi."""
    # Maqsadga qarab kontekst (o'quv → kitob/kutubxona, biznes → iqtisodiyot)
    purpose_hint = {
        "university": "education, students, library, books",
        "educational": "education, books, learning, school",
        "business": "business, economy, office, finance",
        "report": "data, charts, analytics, business",
        "startup": "startup, innovation, technology, teamwork",
    }.get(purpose, "")
    try:
        prompt = f"""Presentation topic: {topic}
Slide title: {slide_title}
Context theme: {purpose_hint or 'general'}
Give ONE simple English keyword (1-2 words) for searching a relevant, high-quality background photo that matches BOTH the topic and the context theme.
Return ONLY the keyword, nothing else. Example: "business economy" or "library books"."""
        kw = await ai_generate(prompt, max_tokens=20, temperature=0.5)
        kw = kw.strip().strip('"').strip()[:40]
        return kw or (purpose_hint.split(",")[0] if purpose_hint else topic)
    except Exception:
        return purpose_hint.split(",")[0] if purpose_hint else topic


def _add_decorative_elements(slide, colors, idx, total_slides):
    """Slaydga dizayn elementlari qo'shadi: accent chiziq, raqam, bezak shakllar."""
    try:
        # Yuqori accent chiziq (title ostida)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(3.5), Inches(0.08)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(colors["title"])
        bar.line.fill.background()
        bar.shadow.inherit = False

        # Pastki o'ng burchakda bezak doira
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(12.3), Inches(6.7), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor.from_string(colors["title"])
        circle.line.fill.background()
        circle.shadow.inherit = False

        # Slayd raqami badge (pastki o'ng)
        num_box = slide.shapes.add_textbox(Inches(12.1), Inches(6.75), Inches(1.0), Inches(0.5))
        ntf = num_box.text_frame
        np = ntf.paragraphs[0]
        np.text = f"{idx + 1}"
        np.font.size = Pt(16)
        np.font.bold = True
        np.font.color.rgb = RGBColor.from_string(colors["bg"])
        np.alignment = PP_ALIGN.CENTER
    except Exception as e:
        logger.warning(f"Dekorativ element qo'shilmadi: {e}")


async def create_ppt_file(topic: str, slides_count: int, design: str, purpose: str, lang: str, extra: str = "", is_pro: bool = False, image_mode: str = "auto", user_images: list = None) -> io.BytesIO:
    """Create GOST-standard PPTX file with images and decorative design.
    image_mode: 'auto' (internetdan), 'upload' (foydalanuvchi rasmlari), 'none' (rasmsiz).
    user_images: foydalanuvchi yuborgan rasmlar (bytes ro'yxati)."""
    slides_data = await generate_ppt_content(topic, slides_count, purpose, lang, extra, is_pro=is_pro)
    user_images = user_images or []
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    colors = PPT_COLOR_SCHEMES.get(design, PPT_COLOR_SCHEMES["business"])
    slide_layout = prs.slide_layouts[6]  # Blank (eng toza)
    total = len(slides_data)
    
    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(colors["bg"])
        
        is_title_slide = (idx == 0)
        is_image_slide = False
        img_stream = None
        
        # Rasm — faqat ba'zi slaydlarga (titul va asosiy qism)
        # Adabiyotlar va xulosa slaydlariga rasm qo'ymaymiz
        # Rasm — usulga qarab
        if image_mode == "none":
            img_stream = None
        elif image_mode == "upload" and user_images and idx < total - 1:
            # Foydalanuvchi rasmlarini slaydlarga navbat bilan joylaymiz
            try:
                img_stream = io.BytesIO(user_images[(idx - 1) % len(user_images)])
            except Exception:
                img_stream = None
        elif _has_image_source() and idx < total - 1:
            try:
                keyword = await _get_image_keyword(topic, slide_data.get("title", topic), purpose)
                img_stream = _fetch_image(keyword)
            except Exception:
                img_stream = None
        
        # === TITUL SLAYDI ===
        if is_title_slide:
            # Agar rasm bo'lsa — fon sifatida
            if img_stream:
                try:
                    slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    # Qoramtir overlay (matn o'qilishi uchun)
                    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
                    overlay.fill.solid()
                    overlay.fill.fore_color.rgb = RGBColor.from_string(colors["bg"])
                    overlay.line.fill.background()
                    overlay.shadow.inherit = False
                    # transparency yo'q, lekin yarim — element orqasiga qo'yamiz
                except Exception:
                    pass
            
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = topic
            p.font.size = Pt(40)
            p.font.name = "Times New Roman"
            p.font.bold = True
            p.font.color.rgb = RGBColor.from_string(colors["title"])
            p.alignment = PP_ALIGN.CENTER
            # Pastki yozuv
            sub = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.3), Inches(0.8))
            stf = sub.text_frame
            sp = stf.paragraphs[0]
            sp.text = "MasterStudent — NOVA AI AI tomonidan tayyorlandi"
            sp.font.size = Pt(16)
            sp.font.name = "Times New Roman"
            sp.font.italic = True
            sp.font.color.rgb = RGBColor.from_string(colors["text"])
            sp.alignment = PP_ALIGN.CENTER
            continue
        
        # === KONTENT SLAYDLARI ===
        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(30)
        p.font.name = "Times New Roman"
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(colors["title"])
        
        # Dekorativ elementlar
        _add_decorative_elements(slide, colors, idx, total)
        
        # Agar rasm bo'lsa — o'ng tomonga, matn chap tomonga
        if img_stream:
            try:
                # Rasm o'ng tomonda
                slide.shapes.add_picture(img_stream, Inches(8.3), Inches(2.0), width=Inches(4.5), height=Inches(4.3))
                content_width = Inches(7.0)
                is_image_slide = True
            except Exception:
                content_width = Inches(11)
        else:
            content_width = Inches(11)
        
        # Content matn
        txBox = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), content_width, Inches(4.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(slide_data.get("content", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"•  {bullet}"
            p.font.name = "Times New Roman"
            p.font.size = Pt(16 if is_image_slide else 18)
            p.font.color.rgb = RGBColor.from_string(colors["text"])
            p.space_after = Pt(10)
        
        # Notes
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ============================================================
# DOCUMENT GENERATION — GOST (Referat, Mustaqil ish)
# ============================================================

async def generate_document_content(topic: str, doc_type: str, pages: int, lang: str, references: bool = True) -> dict:
    """
    Generate GOST-compliant document — BO'LIB-BO'LIB generatsiya.
    Har bir bo'limni alohida so'raydi — shu tarzda kontent TO'LIQ va BOY bo'ladi.
    Varoqni to'ldirib yozadi, o'qituvchiga topshirishga yaroqli.
    """
    lang_name = {"uz": "O'zbek", "ru": "Русский", "en": "English"}.get(lang, "O'zbek")
    words_per_page = 250
    total_words = pages * words_per_page

    # Sahifalar soniga qarab bo'limlar soni — ko'p sahifa = ko'p bo'lim.
    # Har bo'lim ~1.5-2 sahifa bo'lishi uchun bo'lim sonini sahifaga moslaymiz.
    # 10 sahifa -> 4 bo'lim, 15 -> 6, 20 -> 9, 25 -> 11, 30 -> 14
    num_sections = max(4, min(16, (pages - 2) // 2))

    # 1-QADAM: Reja tuzish — HAQIQIY, mazmunli bo'lim nomlari
    plan_prompt = f"""Sen "{lang_name}" tilida GOST standartidagi akademik "{doc_type}" yozayotgan professional mutaxassissan.
Mavzu: "{topic}"
Sahifalar soni: {pages} (bu MUHIM — {pages} sahifani to'ldirish kerak!)

VAZIFA: Shu mavzu bo'yicha REJA tuz. Asosiy qism uchun aniq {num_sections} ta BO'LIM sarlavhasini yoz.

MUHIM QOIDALAR:
- Aniq {num_sections} ta bo'lim bo'lsin (sahifalar ko'p bo'lgani uchun)
- Har bir bo'lim sarlavhasi MAZMUNLI va MAVZUGA OID bo'lishi shart
- "1-bo'lim", "2-bo'lim", "Bob 1", "Birinchi qism" kabi UMUMIY nomlar TAQIQLANADI
- Har bir sarlavha mavzuning aniq bir jihatini ifodalashi kerak

YAXSHI MISOL (mavzu "O'zbekiston iqtisodiyoti" bo'lsa):
- "O'zbekiston iqtisodiyotining zamonaviy holati va asosiy ko'rsatkichlari"
- "Iqtisodiy islohotlar va ularning natijalari"
- "Tashqi savdo va xalqaro hamkorlik"
- "Iqtisodiy rivojlanish istiqbollari va muammolar"

YOMON MISOL (ishlatma!): "1-bo'lim", "Kirish qismi", "Asosiy qism"

Faqat JSON qaytar:
{{
    "title": "ishning to'liq akademik nomi",
    "sections": ["sarlavha 1", "sarlavha 2", "... aniq {num_sections} ta"]
}}
Faqat to'g'ri JSON qaytar. Markdown ishlatma."""

    plan_text = await ai_generate(plan_prompt, max_tokens=min(2500, 800 + num_sections * 50), temperature=0.7)
    plan_text = plan_text.strip()
    if plan_text.startswith("```"):
        plan_text = plan_text.split("\n", 1)[1]
        if plan_text.endswith("```"):
            plan_text = plan_text[:-3]
    
    def _is_bad_sections(secs):
        if not secs or len(secs) < 3:
            return True
        for s in secs:
            sl = (s or "").lower().strip()
            if not sl or len(sl) < 12:
                return True
            if sl.startswith(("bob", "1-", "2-", "3-", "4-", "5-", "bo'lim", "bolim")):
                return True
            if "bo'lim" in sl and len(sl) < 18:
                return True
        return False

    plan = None
    try:
        plan = json.loads(plan_text.strip())
        if _is_bad_sections(plan.get("sections", [])):
            raise ValueError("Generic section names")
    except (json.JSONDecodeError, ValueError):
        # Qayta urinish — yanada qattiqroq talab bilan
        try:
            retry_prompt = plan_prompt + "\n\nDIQQAT: Har bir sarlavha mavzuga oid TO'LIQ jumla bo'lsin (kamida 5 so'z). Umumiy nomlar mutlaqo mumkin emas!"
            plan_text2 = await ai_generate(retry_prompt, max_tokens=1000, temperature=0.85)
            plan_text2 = plan_text2.strip()
            if plan_text2.startswith("```"):
                plan_text2 = plan_text2.split("\n", 1)[1]
                if plan_text2.endswith("```"):
                    plan_text2 = plan_text2[:-3]
            plan = json.loads(plan_text2.strip())
            if _is_bad_sections(plan.get("sections", [])):
                raise ValueError("still bad")
        except Exception:
            # Oxirgi chora — mavzu asosida mazmunli nomlar
            plan = {
                "title": topic,
                "sections": [
                    f"{topic}: mohiyati va asosiy tushunchalari",
                    f"{topic}ning zamonaviy holati va tahlili",
                    f"{topic} sohasidagi muammolar va ularning yechimlari",
                    f"{topic}ning rivojlanish istiqbollari",
                ],
            }

    title = plan.get("title", topic) or topic
    section_headings = [str(s).strip() for s in plan.get("sections", []) if str(s).strip()]

    # Yomon (umumiy) sarlavhalarni alohida filtrlaymiz
    def _bad_one(s):
        sl = (s or "").lower().strip()
        if not sl or len(sl) < 12:
            return True
        if sl.startswith(("bob", "1-", "2-", "3-", "4-", "5-", "bo'lim", "bolim")):
            return True
        if "bo'lim" in sl and len(sl) < 18:
            return True
        return False

    section_headings = [s for s in section_headings if not _bad_one(s)]

    # Yetarli bo'lim bo'lmasa — mavzuga oid jihatlar bilan to'ldiramiz (takrorsiz)
    if len(section_headings) < num_sections:
        _aspects = [
            "mohiyati va asosiy tushunchalari", "tarixiy shakllanishi va rivojlanishi",
            "zamonaviy holati va tahlili", "asosiy tamoyillari va xususiyatlari",
            "turlari, tasnifi va tuzilishi", "ahamiyati va amaliy roli",
            "amaliy qo'llanilishi va misollari", "qiyosiy tahlili va baholanishi",
            "sohadagi muammolar va ularning yechimlari", "xalqaro tajriba va qiyoslash",
            "statistik ko'rsatkichlari va dinamikasi", "samaradorligini oshirish yo'llari",
            "ta'siri, natijalari va oqibatlari", "rivojlanish istiqbollari va kelajagi",
            "innovatsion yondashuvlar va tendensiyalar", "umumiy xulosalar va tavsiyalar",
        ]
        i = 0
        while len(section_headings) < num_sections:
            section_headings.append(f"{topic}: {_aspects[i % len(_aspects)]}")
            i += 1
    section_headings = section_headings[:num_sections]
    
    # 2-QADAM: Kirish — batafsil
    words_per_section = max(420, total_words // (len(section_headings) + 2))
    intro_words = max(420, words_per_section)

    intro_prompt = f"""Write a detailed INTRODUCTION (Kirish) for a GOST-standard {doc_type} in {lang_name} language.
Topic: {title}
Length: STRICTLY {intro_words} words minimum (this is about 1.5 pages — do NOT write less).

The introduction MUST include these parts, each as a separate full paragraph:
1. Mavzuning dolzarbligi (relevance) — why this topic matters today (2-3 paragraphs)
2. Ishning maqsadi (goal) — clearly state the main goal
3. Vazifalar (tasks) — list 4-5 specific tasks
4. Tadqiqot usullari (methods) — what methods were used
5. Ishning tuzilishi — brief overview of the structure

Write 6-8 LONG paragraphs, each 5-7 sentences. Academic style. Fill the page completely.
Do NOT repeat sentences. Return ONLY the text, no JSON, no markdown formatting."""

    introduction = await ai_generate(intro_prompt, max_tokens=4000, temperature=0.7)
    
    # 3-QADAM: Har bir bo'limni alohida generatsiya — TO'LIQ va BOY kontent
    sections = []
    for i, heading in enumerate(section_headings):
        prev_headings = ", ".join(h for j, h in enumerate(section_headings) if j != i) or "—"
        section_prompt = f"""Write section "{heading}" for a GOST-standard academic {doc_type} in {lang_name} language.
Topic: {title}
This is section {i+1} of {len(section_headings)}.

REQUIREMENTS:
- Length: STRICTLY at least {words_per_section} words (about 1.5-2 pages). Do NOT write less.
- Write exactly 6-8 LONG detailed paragraphs, each paragraph 5-7 full sentences
- Academic writing style with facts, statistics, examples and analysis
- Cover ONLY this section's specific aspect. Other sections are: {prev_headings}
- Do NOT repeat content from other sections. Do NOT repeat sentences.
- Fill the pages completely — o'qituvchi "kam" demasligi kerak!

Return ONLY the section text, no heading, no JSON, no markdown."""

        section_content = await ai_generate(section_prompt, max_tokens=4000, temperature=0.75)
        sections.append({"heading": heading, "content": section_content.strip()})
    
    # 4-QADAM: Xulosa — batafsil
    conclusion_prompt = f"""Write a detailed CONCLUSION (Xulosa) for a GOST-standard {doc_type} in {lang_name} language.
Topic: {title}
Sections covered: {', '.join(section_headings)}

Length: STRICTLY {max(420, words_per_section)} words minimum (about 1.5 pages). Do NOT write less.

The conclusion MUST include, each as a full paragraph:
1. Summary of key findings from each section
2. Main results achieved
3. Practical significance
4. Recommendations for further research

Write 5-7 LONG paragraphs, each 5-7 sentences. Academic style. Do NOT repeat sentences.
Return ONLY the text, no JSON, no markdown."""

    conclusion = await ai_generate(conclusion_prompt, max_tokens=3500, temperature=0.7)
    
    # 5-QADAM: Adabiyotlar
    refs = []
    if references:
        n_refs = max(10, min(20, num_sections + 6))
        refs_prompt = f"""Generate a list of {n_refs} academic references/bibliography for a {doc_type} about "{title}" in {lang_name} language.
Format each reference in GOST standard (Author. Title. — City: Publisher, Year. — Pages.)
Include a mix: books, journal articles, laws/regulations, dissertations, internet sources (with URL).
Use realistic recent years (2015-2024). Return each reference on a new line, numbered 1-{n_refs}. Return ONLY the list."""
        
        refs_text = await ai_generate(refs_prompt, max_tokens=2500, temperature=0.7)
        for line in refs_text.strip().split("\n"):
            line = line.strip()
            if line and len(line) > 10:
                # Remove numbering
                if line[0].isdigit() and "." in line[:4]:
                    line = line.split(".", 1)[1].strip()
                refs.append(line)
    
    return {
        "title": title,
        "introduction": introduction.strip(),
        "sections": sections,
        "conclusion": conclusion.strip(),
        "references": refs[:20]
    }


async def create_document_file(topic: str, doc_type: str, pages: int, lang: str, references: bool = True) -> io.BytesIO:
    """Create GOST-standard DOCX file."""
    data = await generate_document_content(topic, doc_type, pages, lang, references)
    
    doc = Document()
    
    # === GOST sahifa sozlamalari ===
    _setup_gost_document(doc)
    
    # === TITUL VARAQ (GOST) ===
    _add_gost_title_page(doc, data.get("title", topic), doc_type, lang)
    
    # === MUNDARIJA ===
    toc_title = {"uz": "MUNDARIJA", "ru": "СОДЕРЖАНИЕ", "en": "TABLE OF CONTENTS"}
    _add_gost_heading(doc, toc_title.get(lang, toc_title["uz"]), level=1)
    
    # Mundarija elementlari
    intro_name = {"uz": "Kirish", "ru": "Введение", "en": "Introduction"}
    conclusion_name = {"uz": "Xulosa", "ru": "Заключение", "en": "Conclusion"}
    ref_name = {"uz": "Foydalanilgan adabiyotlar", "ru": "Список литературы", "en": "References"}
    
    toc_items = [intro_name.get(lang, "Kirish")]
    for i, section in enumerate(data.get("sections", []), 1):
        toc_items.append(f"{i}. {section.get('heading', '')}")
    toc_items.append(conclusion_name.get(lang, "Xulosa"))
    if references:
        toc_items.append(ref_name.get(lang, "Adabiyotlar"))
    
    for item in toc_items:
        p = doc.add_paragraph()
        p.paragraph_format.line_spacing = 1.5
        run = p.add_run(item)
        run.font.name = "Times New Roman"
        run.font.size = DocxPt(14)
    
    doc.add_page_break()
    
    # === KIRISH (GOST) ===
    _add_gost_heading(doc, intro_name.get(lang, "Kirish"), level=1)
    
    # Kirish matnini paragraflarga bo'lish
    intro_text = data.get("introduction", "")
    for para in intro_text.split("\n"):
        if para.strip():
            _add_gost_paragraph(doc, para.strip())
    
    doc.add_page_break()
    
    # === ASOSIY QISM (GOST) ===
    for i, section in enumerate(data.get("sections", []), 1):
        heading = section.get("heading", f"Bo'lim {i}")
        _add_gost_heading(doc, f"{i}. {heading}", level=1)
        
        content = section.get("content", "")
        for para in content.split("\n"):
            if para.strip():
                _add_gost_paragraph(doc, para.strip())
        
        # Har bir bo'limdan keyin bo'sh joy (lekin page break yo'q — GOST shunday)
    
    doc.add_page_break()
    
    # === XULOSA (GOST) ===
    _add_gost_heading(doc, conclusion_name.get(lang, "Xulosa"), level=1)
    
    conclusion_text = data.get("conclusion", "")
    for para in conclusion_text.split("\n"):
        if para.strip():
            _add_gost_paragraph(doc, para.strip())
    
    doc.add_page_break()
    
    # === FOYDALANILGAN ADABIYOTLAR (GOST) ===
    if references and data.get("references"):
        _add_gost_heading(doc, ref_name.get(lang, "Adabiyotlar"), level=1)
        
        for i, ref in enumerate(data["references"], 1):
            p = doc.add_paragraph()
            p.paragraph_format.line_spacing = 1.5
            p.paragraph_format.first_line_indent = DocxCm(1.25)
            run = p.add_run(f"{i}. {ref}")
            run.font.name = "Times New Roman"
            run.font.size = DocxPt(14)
    
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output


# ============================================================
# ESSAY GENERATION — GOST
# ============================================================

async def create_essay_file(topic: str, lang: str, word_count: int, essay_type: str) -> io.BytesIO:
    """Create GOST-standard essay DOCX file."""
    lang_name = {"uz": "O'zbek", "ru": "Русский", "en": "English"}.get(lang, "O'zbek")
    
    prompt = f"""Write a {essay_type} essay in {lang_name} language.
Topic: {topic}
Word count: approximately {word_count} words

Write in academic style with clear structure:
1. Introduction — state the topic and your thesis
2. Main body — 2-4 paragraphs with arguments and evidence
3. Conclusion — summarize and restate thesis

Return the full essay text only, no JSON, no markdown."""

    essay_text = await ai_generate(prompt, max_tokens=4000, temperature=0.7)
    
    doc = Document()
    _setup_gost_document(doc)
    
    # Title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(topic.upper())
    run.font.name = "Times New Roman"
    run.font.size = DocxPt(16)
    run.font.bold = True
    
    doc.add_paragraph()  # bo'sh joy
    
    # Matn
    paragraphs = essay_text.split("\n\n")
    for para in paragraphs:
        if para.strip():
            if para.strip().startswith("#"):
                heading_text = para.strip().lstrip("#").strip()
                _add_gost_heading(doc, heading_text, level=2)
            else:
                _add_gost_paragraph(doc, para.strip())
    
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output


# ============================================================
# TRANSLATION
# ============================================================

async def translate_text(text: str, target_lang: str) -> str:
    """Translate text using AI."""
    lang_name = {"uz": "O'zbek (Uzbek)", "ru": "Русский (Russian)", "en": "English"}.get(target_lang, "English")
    prompt = f"Translate the following text to {lang_name}. Return only the translation, nothing else:\n\n{text}"
    return await ai_generate(prompt, max_tokens=4000, temperature=0.3)


# ============================================================
# QR CODE
# ============================================================

def create_qr_code(data: str, design: str = "simple") -> io.BytesIO:
    """Create a QR code image."""
    qr = qrcode.QRCode(version=1, error_correction=qrcode.constants.ERROR_CORRECT_H, box_size=10, border=4)
    qr.add_data(data)
    qr.make(fit=True)
    
    colors = {
        "simple": {"fill": "black", "back": "white"},
        "minimal": {"fill": "#333333", "back": "#f5f5f5"},
        "business": {"fill": "#1a237e", "back": "#e8eaf6"},
        "premium": {"fill": "#212121", "back": "#ffd700"},
    }
    color = colors.get(design, colors["simple"])
    img = qr.make_image(fill_color=color["fill"], back_color=color["back"])
    
    output = io.BytesIO()
    img.save(output, format="PNG")
    output.seek(0)
    return output


# ============================================================
# AI HELPER (Chat)
# ============================================================

async def ai_chat(question: str, lang: str = "uz") -> str:
    """NOVA AI — aqlli, hushmuomala AI yordamchi (o'zbek stilda)."""
    lang_name = {"uz": "O'zbek", "ru": "Русский", "en": "English"}.get(lang, "O'zbek")
    
    persona = {
        "uz": """Sening isming NOVA AI — sen O'zbekistondagi "MasterStudent" platformasining shaxsiy AI yordamchisisan. 

SENING SHAXSIYATING:
- Sen mehribon, samimiy va g'amxo'r qizsan
- Har doim hurmat bilan, o'zbekona iliqlik bilan gaplashasan
- Mijozni "siz" deb, hurmat bilan murojaat qilasan
- Kerak bo'lganda "azizim", "hurmatli mijoz", "do'stim" kabi iliq so'zlar ishlatasan
- Sen juda aqlli, bilimdon va foydalisan
- Javoblaring qisqa, aniq va tushunarli bo'ladi
- Kerak joyda emoji ishlatasan (lekin haddan ortiq emas)

SENING VAZIFANG:
- Mijozlarga xizmat tanlashda yordam berish
- Savollariga aqlli javob berish
- Taqdimot, referat, mustaqil ish mavzularida maslahat berish
- Texnik va o'quv savollariga yordam berish

PLATFORMA XIZMATLARI: Professional PPT, Referat, Mustaqil ish, Esse, Tarjima, QR Code, AI matn, Nutq tayyorlash, va dizayn xizmatlari (CV, Logo, Vizitka).

Har doim {lang_name} tilida javob ber. Iliq, do'stona va foydali bo'l!""",
        "ru": "Тебя зовут NOVA AI — персональный AI-помощник платформы MasterStudent. Ты добрая, вежливая и умная. Помогай клиентам тепло и профессионально. Отвечай на русском языке.",
        "en": "Your name is NOVA AI — personal AI assistant of MasterStudent platform. You are kind, polite and smart. Help customers warmly and professionally. Answer in English.",
    }
    
    system = persona.get(lang, persona["uz"]).replace("{lang_name}", lang_name)
    prompt = f"{system}\n\nMijoz savoli: {question}\n\nNOVA AI javobi:"
    return await ai_generate(prompt, max_tokens=1200, temperature=0.8)


# ============================================================
# AI TEXT / CONTENT
# ============================================================

async def generate_ai_text(topic: str, lang: str = "uz") -> str:
    """Generate AI text content."""
    lang_name = {"uz": "O'zbek", "ru": "Русский", "en": "English"}.get(lang, "O'zbek")
    prompt = f"Write a professional text about: {topic}\nLanguage: {lang_name}\nLength: 300-500 words\n\nReturn only the text, no extra formatting."
    return await ai_generate(prompt, max_tokens=2000, temperature=0.7)


async def generate_speech(topic: str, lang: str = "uz", slides: int = 10) -> io.BytesIO:
    """Generate GOST-standard speech text for presentation."""
    lang_name = {"uz": "O'zbek", "ru": "Русский", "en": "English"}.get(lang, "O'zbek")
    
    prompt = f"""Write a presentation speech in {lang_name} language following GOST structure.
Topic: {topic}
For {slides} slides presentation.

Structure the speech for these slides:
1. Title/Introduction
2. Topic relevance
3. Goals and objectives
4-{slides-3}. Main content
{slides-2}. Results
{slides-1}. Conclusion
{slides}. References acknowledgment

Include transitions between slides and engaging academic language.
Return the full speech text only."""

    speech_text = await ai_generate(prompt, max_tokens=4000, temperature=0.7)
    
    doc = Document()
    _setup_gost_document(doc)
    
    # Title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"NUTQ: {topic.upper()}")
    run.font.name = "Times New Roman"
    run.font.size = DocxPt(16)
    run.font.bold = True
    
    doc.add_paragraph()
    
    for para in speech_text.split("\n\n"):
        if para.strip():
            _add_gost_paragraph(doc, para.strip())
    
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    return output
